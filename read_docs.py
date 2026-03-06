import PyPDF2
from docx import Document
from pptx import Presentation
import os

# Base path
base_path = r"c:\Users\rohan\OneDrive\Desktop\stress raw 2\project details"

def read_pdf(file_path):
    """Read PDF file and extract text"""
    print(f"\n{'='*80}")
    print(f"Reading: {os.path.basename(file_path)}")
    print('='*80)
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page_num, page in enumerate(pdf_reader.pages, 1):
                page_text = page.extract_text()
                text += f"\n--- Page {page_num} ---\n{page_text}"
            print(text)
    except Exception as e:
        print(f"Error reading PDF: {e}")

def read_docx(file_path):
    """Read Word document and extract text"""
    print(f"\n{'='*80}")
    print(f"Reading: {os.path.basename(file_path)}")
    print('='*80)
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                print(para.text)
    except Exception as e:
        print(f"Error reading DOCX: {e}")

def read_pptx(file_path):
    """Read PowerPoint presentation and extract text"""
    print(f"\n{'='*80}")
    print(f"Reading: {os.path.basename(file_path)}")
    print('='*80)
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            print(f"\n--- Slide {slide_num} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        print(shape.text)
    except Exception as e:
        print(f"Error reading PPTX: {e}")

# Read all files
files = [
    ("project details.pdf", read_pdf),
    ("project details1.pdf", read_pdf),
    ("literature survey.pdf", read_pdf),
    ("project flow details.docx", read_docx),
    ("healthcare n lifestyle phase -02 (1) (1).pptx", read_pptx),
]

for filename, reader_func in files:
    file_path = os.path.join(base_path, filename)
    if os.path.exists(file_path):
        reader_func(file_path)
    else:
        print(f"\nFile not found: {filename}")
